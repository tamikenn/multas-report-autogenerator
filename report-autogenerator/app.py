import streamlit as st
import pandas as pd
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
import tempfile
import os

st.title('Excel発表資料自動生成ツール')

uploaded_file = st.file_uploader('Excelファイルをアップロードしてください', type=['xlsx'])

if uploaded_file:
    # Excelファイルの読み込み
    df = pd.read_excel(uploaded_file)
    st.write('アップロードされたデータ:')
    st.dataframe(df)

    # PPT生成
    pptx_path = os.path.join(tempfile.gettempdir(), 'output.pptx')
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = '自動生成スライド'
    # 1行目をテキストとして追加
    text = '\n'.join([str(row) for row in df.iloc[0]])
    slide.shapes.add_textbox(100, 100, 500, 300).text = text
    prs.save(pptx_path)
    with open(pptx_path, 'rb') as f:
        st.download_button('PPTXダウンロード', f, file_name='output.pptx')

    # PDF生成
    pdf_path = os.path.join(tempfile.gettempdir(), 'output.pdf')
    c = canvas.Canvas(pdf_path, pagesize=A4)
    c.drawString(100, 800, '自動生成PDF')
    for i, row in enumerate(df.values.tolist()[:10]):
        c.drawString(100, 780 - i*20, str(row))
    c.save()
    with open(pdf_path, 'rb') as f:
        st.download_button('PDFダウンロード', f, file_name='output.pdf')

    # rawデータ（CSV）
    csv = df.to_csv(index=False).encode('utf-8')
    st.download_button('rawデータ(CSV)ダウンロード', csv, file_name='output.csv')
